#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPTX 화면정의서에서 이미지를 추출하는 스크립트

이미지 추출 워크플로우:
1. PPTX 파일에서 모든 이미지 추출
2. output/images/ 폴더에 저장
3. image_manifest.json 생성 (메타데이터)

보안 기능:
- .gitignore 자동 생성
- 민감 정보 경고 메시지
- cleanup 함수 제공

사용법:
    py extract_images.py "화면정의서.pptx" --output "output"
"""

import argparse
import json
import sys
import shutil
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def check_win32com_available() -> bool:
    """win32com 사용 가능 여부 확인 (pywin32 + PowerPoint 설치 필요)"""
    try:
        import win32com.client
        import pythoncom
        return True
    except ImportError:
        return False


def print_security_warning():
    """보안 경고 메시지 출력"""
    print()
    print("=" * 60)
    print("  [보안 알림]")
    print("=" * 60)
    print("  화면정의서에서 이미지를 추출합니다.")
    print("  - 민감한 정보(환자 데이터, 개인정보 등)가 포함된 경우 주의하세요.")
    print("  - 추출된 이미지는 output/images/ 폴더에 저장됩니다.")
    print("  - --cleanup 옵션으로 TC 생성 후 자동 삭제할 수 있습니다.")
    print("=" * 60)
    print()


def create_gitignore(output_dir: Path):
    """output 폴더에 .gitignore 생성

    추출된 이미지와 임시 JSON 파일이 버전 관리에 포함되지 않도록 합니다.
    """
    gitignore_path = output_dir / ".gitignore"
    gitignore_content = """# 테스트케이스 생성기 임시 파일
# 이 폴더의 이미지는 화면정의서에서 추출된 것으로
# 민감한 정보가 포함될 수 있습니다.

images/
image_manifest.json
image_analysis.json
*.png
*.jpg
*.jpeg
*.gif
*.bmp
*.tiff

# 분석 결과는 유지 (필요시 주석 해제)
# !image_analysis.json
"""
    with open(gitignore_path, "w", encoding="utf-8") as f:
        f.write(gitignore_content)

    print(f"  .gitignore 생성됨: {gitignore_path}")


def cleanup_images(output_dir: str, keep_manifest: bool = False):
    """TC 생성 완료 후 이미지 파일 삭제

    Args:
        output_dir: 출력 폴더 경로
        keep_manifest: manifest 파일 유지 여부 (기본: False)
    """
    output_path = Path(output_dir)
    images_dir = output_path / "images"

    deleted_count = 0

    # 이미지 폴더 삭제
    if images_dir.exists():
        file_count = len(list(images_dir.glob("*")))
        shutil.rmtree(images_dir)
        deleted_count += file_count
        print(f"  이미지 폴더 삭제됨: {images_dir} ({file_count}개 파일)")

    # manifest 삭제 (옵션)
    if not keep_manifest:
        manifest_path = output_path / "image_manifest.json"
        if manifest_path.exists():
            manifest_path.unlink()
            print(f"  manifest 삭제됨: {manifest_path}")

    # 분석 결과 삭제 (선택적)
    analysis_path = output_path / "image_analysis.json"
    if analysis_path.exists() and not keep_manifest:
        analysis_path.unlink()
        print(f"  분석 결과 삭제됨: {analysis_path}")

    print(f"  총 {deleted_count}개 파일 삭제 완료")


def get_image_extension(content_type: str) -> str:
    """컨텐츠 타입에서 확장자 추출"""
    extension_map = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/bmp": ".bmp",
        "image/tiff": ".tiff",
        "image/x-wmf": ".wmf",
        "image/x-emf": ".emf",
    }
    return extension_map.get(content_type, ".png")


def export_slides_fullpage(pptx_path: Path, output_dir: Path,
                           width: int = 1920, height: int = 1080) -> list:
    """win32com COM으로 슬라이드를 통째로 PNG로 내보내기

    PowerPoint Application COM 객체를 사용하여 각 슬라이드를 지정 해상도로 내보냅니다.

    Args:
        pptx_path: PPTX 파일 절대 경로
        output_dir: 출력 폴더 경로
        width: 내보내기 가로 해상도 (px)
        height: 내보내기 세로 해상도 (px)

    Returns:
        fullpage 이미지 정보 리스트
    """
    import pythoncom
    import win32com.client

    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    fullpage_images = []
    ppt = None
    presentation = None

    try:
        pythoncom.CoInitialize()

        ppt = win32com.client.Dispatch("PowerPoint.Application")
        # PowerPoint 창을 표시하지 않음 (백그라운드 실행)
        ppt.DisplayAlerts = 0

        # PPTX 파일 열기 (ReadOnly=True, WithWindow=False)
        abs_path = str(pptx_path.resolve())
        presentation = ppt.Presentations.Open(
            abs_path,
            ReadOnly=True,
            Untitled=False,
            WithWindow=False
        )

        slide_count = presentation.Slides.Count
        print(f"  [Fullpage] {slide_count}개 슬라이드 내보내기 중...")

        for i in range(1, slide_count + 1):
            filename = f"slide_{i:02d}_fullpage.png"
            export_path = str((images_dir / filename).resolve())

            try:
                presentation.Slides(i).Export(export_path, "PNG", width, height)

                image_info = {
                    "filename": filename,
                    "path": export_path,
                    "slide_number": i,
                    "image_type": "fullpage",
                    "content_type": "image/png",
                    "size": {
                        "width_px": width,
                        "height_px": height
                    }
                }
                fullpage_images.append(image_info)

            except Exception as e:
                print(f"  [Fullpage] 슬라이드 {i} 내보내기 실패: {e}")

        print(f"  [Fullpage] {len(fullpage_images)}/{slide_count}개 내보내기 완료")

    except Exception as e:
        print(f"  [Fullpage] PowerPoint COM 오류: {e}")
        print(f"  [Fullpage] 개별 이미지 추출만 사용됩니다.")

    finally:
        try:
            if presentation:
                presentation.Close()
        except Exception:
            pass
        try:
            if ppt:
                ppt.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass

    return fullpage_images


def extract_images_from_pptx(pptx_path: Path, output_dir: Path) -> dict:
    """PPTX에서 모든 이미지를 추출하여 저장

    Args:
        pptx_path: PPTX 파일 경로
        output_dir: 출력 폴더 경로

    Returns:
        image_manifest: 추출된 이미지 목록과 메타데이터
    """
    prs = Presentation(str(pptx_path))

    # 이미지 출력 폴더 생성
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    # manifest 초기화
    manifest = {
        "pptx_file": pptx_path.name,
        "pptx_path": str(pptx_path.resolve()),
        "extraction_date": datetime.now().isoformat(),
        "output_dir": str(output_dir.resolve()),
        "total_images": 0,
        "images": []
    }

    image_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_image_index = 0

        for shape in slide.shapes:
            # 이미지 shape 확인
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    content_type = image.content_type
                    extension = get_image_extension(content_type)

                    # 파일명 생성
                    filename = f"slide_{slide_num:02d}_image_{slide_image_index:02d}{extension}"
                    image_path = images_dir / filename

                    # 이미지 저장
                    with open(image_path, "wb") as f:
                        f.write(image.blob)

                    # 이미지 크기 정보 (shape의 크기)
                    width = shape.width.inches if hasattr(shape.width, 'inches') else 0
                    height = shape.height.inches if hasattr(shape.height, 'inches') else 0

                    # 위치 정보
                    left = shape.left.inches if hasattr(shape.left, 'inches') else 0
                    top = shape.top.inches if hasattr(shape.top, 'inches') else 0

                    # manifest에 추가
                    image_info = {
                        "filename": filename,
                        "path": str(image_path.resolve()),
                        "slide_number": slide_num,
                        "image_index": slide_image_index,
                        "content_type": content_type,
                        "size": {
                            "width_inches": round(width, 2),
                            "height_inches": round(height, 2)
                        },
                        "position": {
                            "left_inches": round(left, 2),
                            "top_inches": round(top, 2)
                        },
                        "analysis": None  # Claude Code 분석 후 채워짐
                    }
                    manifest["images"].append(image_info)

                    image_count += 1
                    slide_image_index += 1

                except Exception as e:
                    print(f"  [경고] 슬라이드 {slide_num} 이미지 추출 실패: {e}")

    manifest["total_images"] = image_count

    return manifest


def save_manifest(manifest: dict, output_dir: Path):
    """manifest 파일 저장"""
    manifest_path = output_dir / "image_manifest.json"

    with open(manifest_path, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)

    return manifest_path


def extract_and_save_images(pptx_path: str, output_dir: str = None,
                            no_fullpage: bool = False) -> dict:
    """PPTX에서 이미지 추출 및 저장 (메인 함수)

    Args:
        pptx_path: PPTX 파일 경로
        output_dir: 출력 폴더 경로 (기본: ./output)
        no_fullpage: True이면 fullpage 캡처 건너뜀

    Returns:
        manifest: 추출 결과 정보
    """
    pptx_path = Path(pptx_path).resolve()

    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX 파일을 찾을 수 없습니다: {pptx_path}")

    if not pptx_path.suffix.lower() == ".pptx":
        raise ValueError(f"PPTX 파일이 아닙니다: {pptx_path}")

    # 출력 폴더 설정
    if output_dir:
        output_path = Path(output_dir).resolve()
    else:
        output_path = pptx_path.parent / "output"

    output_path.mkdir(parents=True, exist_ok=True)

    # .gitignore 생성
    create_gitignore(output_path)

    # 개별 이미지 추출
    print(f"  PPTX 파일: {pptx_path}")
    manifest = extract_images_from_pptx(pptx_path, output_path)

    # Fullpage 슬라이드 캡처
    if no_fullpage:
        manifest["fullpage_images"] = []
        manifest["total_fullpage_images"] = 0
        manifest["fullpage_export_method"] = "skipped"
        manifest["fullpage_resolution"] = None
    elif check_win32com_available():
        # config에서 해상도 가져오기
        try:
            sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
            from config import FULLPAGE_WIDTH, FULLPAGE_HEIGHT
        except ImportError:
            FULLPAGE_WIDTH, FULLPAGE_HEIGHT = 1920, 1080

        print(f"  [Fullpage] win32com 사용 가능 ({FULLPAGE_WIDTH}x{FULLPAGE_HEIGHT})")
        fullpage_images = export_slides_fullpage(
            pptx_path, output_path, FULLPAGE_WIDTH, FULLPAGE_HEIGHT
        )
        manifest["fullpage_images"] = fullpage_images
        manifest["total_fullpage_images"] = len(fullpage_images)
        manifest["fullpage_export_method"] = "win32com"
        manifest["fullpage_resolution"] = {
            "width": FULLPAGE_WIDTH,
            "height": FULLPAGE_HEIGHT
        }
    else:
        print("  [Fullpage] win32com 미설치 - fullpage 캡처 건너뜀")
        print("  [Fullpage] 설치: pip install pywin32")
        manifest["fullpage_images"] = []
        manifest["total_fullpage_images"] = 0
        manifest["fullpage_export_method"] = "unavailable"
        manifest["fullpage_resolution"] = None

    # manifest 저장
    manifest_path = save_manifest(manifest, output_path)
    print(f"  manifest 생성됨: {manifest_path}")

    return manifest


def main():
    """메인 함수"""
    parser = argparse.ArgumentParser(
        description="PPTX 화면정의서에서 이미지를 추출합니다.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
예시:
  py extract_images.py "화면정의서.pptx"
  py extract_images.py "화면정의서.pptx" --output "output"
  py extract_images.py --cleanup "output"
        """
    )

    parser.add_argument(
        "pptx_file",
        nargs="?",
        help="입력 PPTX 파일 경로"
    )

    parser.add_argument(
        "--output", "-o",
        dest="output_dir",
        default=None,
        help="출력 폴더 경로 (기본값: ./output)"
    )

    parser.add_argument(
        "--cleanup", "-c",
        dest="cleanup_dir",
        metavar="DIR",
        help="지정된 폴더의 추출된 이미지 삭제"
    )

    parser.add_argument(
        "--keep-manifest",
        action="store_true",
        help="cleanup 시 manifest 파일 유지"
    )

    parser.add_argument(
        "--quiet", "-q",
        action="store_true",
        help="보안 경고 메시지 생략"
    )

    parser.add_argument(
        "--no-fullpage",
        action="store_true",
        help="슬라이드 전체 캡처(fullpage) 건너뛰기"
    )

    args = parser.parse_args()

    # cleanup 모드
    if args.cleanup_dir:
        print(f"[정리] 추출된 이미지 삭제 중...")
        cleanup_images(args.cleanup_dir, args.keep_manifest)
        print("[정리] 완료")
        return 0

    # 추출 모드
    if not args.pptx_file:
        parser.print_help()
        return 1

    if not args.quiet:
        print_security_warning()

    print("[추출] 이미지 추출 시작...")

    try:
        manifest = extract_and_save_images(
            pptx_path=args.pptx_file,
            output_dir=args.output_dir,
            no_fullpage=args.no_fullpage
        )

        print()
        print("-" * 60)
        print("  추출 결과 요약")
        print("-" * 60)
        print(f"  입력 파일      : {manifest['pptx_file']}")
        print(f"  개별 이미지    : {manifest['total_images']}개")
        fullpage_count = manifest.get('total_fullpage_images', 0)
        fullpage_method = manifest.get('fullpage_export_method', 'unavailable')
        print(f"  Fullpage 이미지: {fullpage_count}개 ({fullpage_method})")
        print(f"  출력 폴더      : {manifest['output_dir']}")
        print("-" * 60)
        print()
        print("[추출] 완료!")

        return 0

    except FileNotFoundError as e:
        print(f"[오류] 파일을 찾을 수 없습니다: {e}")
        return 1
    except ValueError as e:
        print(f"[오류] 잘못된 입력: {e}")
        return 1
    except Exception as e:
        print(f"[오류] 예상치 못한 오류: {e}")
        import traceback
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
