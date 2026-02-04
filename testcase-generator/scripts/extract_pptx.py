#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPTX 화면정의서에서 컴포넌트 정보를 추출하는 스크립트
"""

import json
import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def extract_table_data(table):
    """테이블에서 데이터 추출"""
    rows = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            text = cell.text.strip() if cell.text else ""
            row_data.append(text)
        rows.append(row_data)
    return rows


def parse_header_table(rows):
    """헤더 테이블 파싱 (프로젝트명, 화면ID 등)"""
    header_info = {}

    for row in rows:
        for i, cell in enumerate(row):
            cell_lower = cell.lower() if cell else ""

            # 키-값 쌍 찾기
            if "프로젝트" in cell or "project" in cell_lower:
                if i + 1 < len(row):
                    header_info["project_name"] = row[i + 1]
            elif "화면" in cell and "id" in cell.lower():
                if i + 1 < len(row):
                    header_info["screen_id"] = row[i + 1]
            elif "제목" in cell or "title" in cell_lower:
                if i + 1 < len(row):
                    header_info["title"] = row[i + 1]
            elif "문서번호" in cell:
                if i + 1 < len(row):
                    header_info["doc_number"] = row[i + 1]
            elif "작성일자" in cell:
                if i + 1 < len(row):
                    header_info["date"] = row[i + 1]
            elif "버전" in cell or "version" in cell_lower:
                if i + 1 < len(row):
                    header_info["version"] = row[i + 1]

    return header_info


def parse_component_table(rows):
    """컴포넌트 테이블 파싱 (No, Component, Description)"""
    components = []

    # 헤더 행 찾기
    header_row_idx = -1
    for i, row in enumerate(rows):
        row_lower = [cell.lower() for cell in row]
        if "no" in row_lower and ("component" in row_lower or "description" in row_lower):
            header_row_idx = i
            break

    if header_row_idx == -1:
        return components

    # 컬럼 인덱스 매핑
    header = rows[header_row_idx]
    col_map = {}
    for i, cell in enumerate(header):
        cell_lower = cell.lower()
        if cell_lower == "no":
            col_map["no"] = i
        elif "component" in cell_lower:
            col_map["component"] = i
        elif "description" in cell_lower:
            col_map["description"] = i

    # 데이터 행 파싱
    for row in rows[header_row_idx + 1:]:
        if len(row) <= max(col_map.values(), default=0):
            continue

        component = {}
        if "no" in col_map:
            no_val = row[col_map["no"]].strip()
            if no_val and no_val.isdigit():
                component["no"] = int(no_val)
            else:
                continue  # 번호가 없으면 건너뛰기

        if "component" in col_map:
            component["component"] = row[col_map["component"]].strip()

        if "description" in col_map:
            component["description"] = row[col_map["description"]].strip()

        if component.get("component"):
            components.append(component)

    return components


def extract_text_from_shape(shape):
    """Shape에서 텍스트 추출"""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""


def extract_slide_info(slide, slide_number):
    """슬라이드에서 정보 추출"""
    slide_info = {
        "slide_number": slide_number,
        "header": {},
        "components": [],
        "section_title": "",
        "raw_text": []
    }

    tables = []
    text_content = []

    for shape in slide.shapes:
        # 테이블 처리
        if shape.has_table:
            table_data = extract_table_data(shape.table)
            tables.append(table_data)

        # 텍스트 처리
        text = extract_text_from_shape(shape)
        if text:
            text_content.append(text)
            slide_info["raw_text"].append(text)

    # 테이블 분석
    for table in tables:
        # 헤더 테이블인지 확인
        first_row_text = " ".join(table[0]) if table else ""
        if "프로젝트" in first_row_text or "화면" in first_row_text:
            slide_info["header"] = parse_header_table(table)

        # 컴포넌트 테이블인지 확인
        for row in table[:2]:  # 처음 2행만 확인
            row_text = " ".join(row).lower()
            if "no" in row_text and "component" in row_text:
                components = parse_component_table(table)
                slide_info["components"].extend(components)
                break

    # 섹션 제목 추출 (예: "01 메인 Layout 및 Tool 설명")
    for text in text_content:
        if re.match(r"^\d{2}\s+", text):
            slide_info["section_title"] = text
            break

    return slide_info


def extract_pptx(pptx_path):
    """PPTX 파일에서 전체 정보 추출"""
    prs = Presentation(pptx_path)

    result = {
        "file_path": str(pptx_path),
        "total_slides": len(prs.slides),
        "slides": []
    }

    for i, slide in enumerate(prs.slides, 1):
        slide_info = extract_slide_info(slide, i)
        result["slides"].append(slide_info)

    # 프로젝트 정보 집계
    for slide in result["slides"]:
        if slide["header"]:
            result["project_info"] = slide["header"]
            break

    # 전체 컴포넌트 집계
    all_components = []
    for slide in result["slides"]:
        for comp in slide["components"]:
            comp["slide_number"] = slide["slide_number"]
            comp["section"] = slide["section_title"]
            all_components.append(comp)

    result["all_components"] = all_components

    return result


def main():
    if len(sys.argv) < 2:
        print("Usage: python extract_pptx.py <pptx_file> [output_json]")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) > 2 else None

    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    result = extract_pptx(pptx_path)

    output_json = json.dumps(result, ensure_ascii=False, indent=2)

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(output_json)
        print(f"Output saved to: {output_path}")
    else:
        print(output_json)

    return result


if __name__ == "__main__":
    main()
